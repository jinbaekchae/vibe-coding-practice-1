"""
교육 제안서 2안 PPTX — 바이브코딩 중심 구조
슬라이드 크기: 10.00in × 5.62in
색상: 기업소개서 테마 (다크 #1E1E1E / 회색 #666666 / 카드 #FCFCFC 등)
장식 최소: 상단 바 없음, 언더라인 없음, 오렌지 채운 박스 없음
총 19장 슬라이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 색상 ────────────────────────────────────────────────────────
DARK = RGBColor(0x1E, 0x1E, 0x1E)   # 제목·본문
WH   = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색 배경
GRAY = RGBColor(0x66, 0x66, 0x66)   # 섹션 번호·부제
CARD = RGBColor(0xFC, 0xFC, 0xFC)   # 카드 배경
HEAD = RGBColor(0xE9, 0xE9, 0xE9)   # 표 헤더
DATA = RGBColor(0xF5, 0xF5, 0xF5)   # 표 데이터
OR   = RGBColor(0xFF, 0x6C, 0x00)   # 오렌지 accent

FONT = "Noto Sans KR"
W = Inches(10.00)
H = Inches(5.62)

# ── 그림자 제거 ────────────────────────────────────────────────
def no_shadow(shape):
    """도형의 그림자(그리고 모든 효과) 제거 — effectLst를 빈 태그로 교체"""
    sp_pr = shape._element.spPr
    for ef in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(ef)
    etree.SubElement(sp_pr, qn('a:effectLst'))

# ── 기본 헬퍼 ──────────────────────────────────────────────────
def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def tb(slide, text, x, y, w, h,
       size=10, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    s = slide.shapes.add_textbox(x, y, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if '\n' in text:
        for para in tf.paragraphs:
            para.line_spacing = 2.0
    no_shadow(s)
    return s

def tb_lines(slide, lines, x, y, w, h,
             size=10, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    """여러 줄 텍스트를 하나의 텍스트박스에 담는다.
    2줄 이상이면 줄간격 2.0, 1줄이면 1.0."""
    s = slide.shapes.add_textbox(x, y, w, h)
    tf = s.text_frame
    tf.word_wrap = True
    multi = len(lines) > 1
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if multi:
            p.line_spacing = 2.0
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    no_shadow(s)
    return s

def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width     = lw
    else:
        s.line.fill.background()
    no_shadow(s)
    return s

# ── 콘텐츠 슬라이드 헤더 ───────────────────────────────────────
def content_page(slide, section_num, title, subtitle=""):
    """섹션 번호 (좌상단 작은 회색) + 제목 (큰 볼드) + 부제"""
    tb(slide, section_num,
       Inches(0.58), Inches(0.18), Inches(5.0), Inches(0.3),
       size=11, color=GRAY)
    tb(slide, title,
       Inches(0.56), Inches(0.42), Inches(8.89), Inches(0.63),
       size=22, bold=True, color=OR)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.56), Inches(0.87), Inches(8.89), Inches(0.3),
           size=10, color=GRAY)

# ── S01 표지 ──────────────────────────────────────────────────
def s01(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=DARK)

    tb(sl, "AI 역량 강화 교육 제안서",
       Inches(0.56), Inches(0.85), Inches(8.89), Inches(1.0),
       size=30, bold=True, color=WH)

    tb(sl, "주식회사 워카 (Woka)  ·  2026년 2월 24일",
       Inches(0.56), Inches(1.9), Inches(8.89), Inches(0.32),
       size=11, color=GRAY)

    rows = [
        ("제안사", "주식회사 워카 (Woka)"),
        ("수신사", "컨트로맥스"),
        ("작성일", "2026년 2월 24일"),
        ("담당자", "교육팀 채진백 팀장  |  010-2326-4348  |  jinbaek@woka.kr"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(2.38) + i * Inches(0.58)
        tb(sl, k, Inches(0.56), y, Inches(1.3), Inches(0.42),
           size=11, color=GRAY)
        tb(sl, v, Inches(1.96), y, Inches(7.5), Inches(0.42),
           size=11, bold=True, color=WH)

    tb(sl, "WOKA",
       Inches(8.4), Inches(5.2), Inches(1.4), Inches(0.3),
       size=14, bold=True, color=OR, align=PP_ALIGN.RIGHT)

# ── S02 목차 ──────────────────────────────────────────────────
def s02(prs):
    sl = blank(prs)
    content_page(sl, "목차", "목  차")

    items = [
        ("01", "교육 의뢰 배경"),
        ("02", "설문 분석 요약"),
        ("03", "교육 과정 제안"),
        ("04", "AI 역량 평가 및 인사평가 연계 체계"),
        ("05", "실행 로드맵"),
        ("06", "워카 소개 & 레퍼런스"),
        ("07", "담당자 연락처"),
    ]
    for i, (num, title) in enumerate(items):
        y = Inches(1.22) + i * Inches(0.52)
        tb(sl, num, Inches(0.56), y, Inches(0.5), Inches(0.38),
           size=11, bold=True, color=OR)
        tb(sl, title, Inches(1.18), y, Inches(8.27), Inches(0.38),
           size=13, color=DARK)

# ── S03 교육 의뢰 배경 ────────────────────────────────────────
def s03(prs):
    sl = blank(prs)
    content_page(
        sl, "1. 교육 의뢰 배경", "교육 의뢰 배경",
        "컨트로맥스는 정밀 서보 액추에이터 제조 분야 소수 정예 기업입니다. 1인 다역 구조에서 AI 활용 생산성 향상이 필수 과제이며, 경영진도 적극 추진 중입니다."
    )

    stats = [
        ("80%",   "교육 참여 의향"),
        ("55%",   "AI 정기 활용률"),
        ("85%",   "업무 시간 단축"),
        ("55%",   "입문·초급 인력"),
    ]
    cw = Inches(2.02)
    gap = Inches(0.27)
    for i, (val, label) in enumerate(stats):
        x = Inches(0.56) + i * (cw + gap)
        y = Inches(1.19)
        rect(sl, x, y, cw, Inches(0.96), fill=CARD)
        tb(sl, val, x, y + Inches(0.06), cw, Inches(0.52),
           size=26, bold=True, color=OR, align=PP_ALIGN.CENTER)
        tb(sl, label, x, y + Inches(0.6), cw, Inches(0.3),
           size=9, color=DARK, align=PP_ALIGN.CENTER)

    tb(sl, "클라이언트 교육 기대사항",
       Inches(0.56), Inches(2.28), Inches(8.89), Inches(0.3),
       size=11, bold=True, color=DARK)

    headers = ["기대사항", "내용", "진행 상황"]
    items_t = [
        ("AI 활용능력 진단",   "초심자 / 중급자 / 고급 수준 진단",         "설문 완료 ✓"),
        ("맞춤형 교육 서비스", "진단 결과별 수준·직무 맞춤 교육",           "과정 A / B / C / D"),
        ("AI 활용능력 평가",   "교육 전·후 성과·수준 향상 측정",           "사전·사후 평가"),
        ("인사평가 반영",      "평가 결과 인사평가 공식 반영",              "레벨업 기준 협의"),
    ]
    cw2 = [Inches(2.5), Inches(3.5), Inches(2.88)]
    xs2 = [Inches(0.56), Inches(3.06), Inches(6.56)]
    rh = Inches(0.32)

    for j, (hdr, xp) in enumerate(zip(headers, xs2)):
        rect(sl, xp, Inches(2.66), cw2[j], rh, fill=HEAD)
        tb(sl, hdr, xp + Inches(0.08), Inches(2.68),
           cw2[j] - Inches(0.16), Inches(0.26),
           size=9, bold=True, color=DARK)
    for i, (c1, c2, c3) in enumerate(items_t):
        ry = Inches(2.98) + i * rh
        for j, (txt, xp) in enumerate(zip([c1, c2, c3], xs2)):
            rect(sl, xp, ry, cw2[j], rh, fill=DATA)
            tb(sl, txt, xp + Inches(0.08), ry + Inches(0.04),
               cw2[j] - Inches(0.16), Inches(0.24),
               size=9, color=DARK)

# ── S04 설문 핵심 수치 ────────────────────────────────────────
def s04(prs):
    sl = blank(prs)
    content_page(sl, "2. 설문 분석 요약", "설문 분석 요약 — 핵심 수치",
                 "응답 인원: 33명  |  조사 일시: 2026년 2월")

    data = [
        ["항목",                          "주요 수치",         "시사점"],
        ["AI 정기 활용률 (주 3회 이상)",  "55%",              "주요 업무 도구로 이미 자리 잡음"],
        ["AI 활용 효과 경험",             "70%",              "'상당한 개선' 또는 '혁신적 변화' 경험"],
        ["주요 도구 (텍스트 AI)",         "90%",              "ChatGPT · Claude · Gemini"],
        ["입문·초급 합산 (Lv1~2)",        "55%",              "맞춤형 기초 교육 필요 인력"],
        ["유료 계정 미보유",               "50%",              "무료 버전만 사용"],
        ["교육 참여 의향 (긍정)",          "80%",              "매우 그렇다 40% + 그렇다 40%"],
        ["업무 시간 단축 경험",            "85%",              "10% 이상 단축 응답자 합산"],
        ["선호 교육 형태",                "월 1회, 4시간 이내", "정기 실습 위주"],
    ]
    cw = [Inches(3.8), Inches(1.7), Inches(3.28)]
    xs = [Inches(0.56), Inches(4.45), Inches(6.24)]
    rh = Inches(0.38)

    for r, row in enumerate(data):
        ry = Inches(1.19) + r * rh
        for j, (cell, xp) in enumerate(zip(row, xs)):
            bg = HEAD if r == 0 else DATA
            rect(sl, xp, ry, cw[j], rh, fill=bg)
            tb(sl, cell, xp + Inches(0.08), ry + Inches(0.06),
               cw[j] - Inches(0.16), rh - Inches(0.12),
               size=10, bold=(r == 0), color=DARK)

# ── S05 업무 영역 & 장벽 ──────────────────────────────────────
def s05(prs):
    sl = blank(prs)
    content_page(sl, "2. 설문 분석 요약", "설문 분석 — AI 도움 업무 영역 & 미활용 장벽")

    tb(sl, "AI가 도움을 주는 주요 업무 영역",
       Inches(0.56), Inches(1.19), Inches(4.1), Inches(0.28),
       size=11, bold=True, color=DARK)
    areas = [
        ("문서 작성 (보고서·이메일)", 70),
        ("아이디어·기획안 초안",     50),
        ("단순 업무보조 (엑셀 수식 등)", 45),
        ("요약·분석",               30),
        ("R&D·코드 생성",           30),
    ]
    bmax = Inches(3.8)
    for i, (label, val) in enumerate(areas):
        y = Inches(1.55) + i * Inches(0.66)
        tb(sl, label, Inches(0.56), y, Inches(3.8), Inches(0.22),
           size=10, color=DARK)
        rect(sl, Inches(0.56), y + Inches(0.25), bmax, Inches(0.17), fill=DATA)
        rect(sl, Inches(0.56), y + Inches(0.25), bmax * val / 100, Inches(0.17), fill=OR)
        tb(sl, f"{val}%", Inches(4.44), y + Inches(0.22),
           Inches(0.45), Inches(0.22), size=10, bold=True, color=OR)

    rect(sl, Inches(4.95), Inches(1.19), Inches(0.005), Inches(4.0), fill=HEAD)

    tb(sl, "미활용 주요 장벽",
       Inches(5.1), Inches(1.19), Inches(4.3), Inches(0.28),
       size=11, bold=True, color=DARK)
    barriers = [
        ("보안 우려 (사내 정보 유출)",        35),
        ("구체적 사용법·프롬프트 미숙",       20),
        ("결과물 신뢰성 부족",               20),
    ]
    bmax2 = Inches(4.0)
    for i, (label, val) in enumerate(barriers):
        y = Inches(1.55) + i * Inches(0.66)
        tb(sl, label, Inches(5.1), y, Inches(4.0), Inches(0.22),
           size=10, color=DARK)
        rect(sl, Inches(5.1), y + Inches(0.25), bmax2, Inches(0.17), fill=DATA)
        rect(sl, Inches(5.1), y + Inches(0.25), bmax2 * val / 100, Inches(0.17), fill=OR)
        tb(sl, f"{val}%", Inches(9.18), y + Inches(0.22),
           Inches(0.45), Inches(0.22), size=10, bold=True, color=OR)

    tb(sl,
       "보안 우려(35%) → 사내 AI 가이드라인 수립  |  프롬프트 미숙(20%) → 과정 A 직접 해결  |  신뢰성 부족(20%) → 과정 A·D 할루시네이션 검증",
       Inches(0.56), Inches(5.2), Inches(8.89), Inches(0.3),
       size=9, color=GRAY)

# ── S06 교육 니즈 3가지 ───────────────────────────────────────
def s06(prs):
    sl = blank(prs)
    content_page(sl, "2. 설문 분석 요약", "도출된 교육 니즈 3가지")

    needs = [
        (
            "니즈 1",
            "AI 기초 리터러시 및 프롬프트 실습",
            "입문·초급 합산 55%  |  사용법 미숙 20%  |  교육 참여 의향 80%\n"
            "선호 \"월 1회·4시간 이내\" 일치  →  과정 A로 직접 해결"
        ),
        (
            "니즈 2",
            "바이브코딩 기반 업무 자동화",
            "문서작성·단순업무보조 합산 115%  |  직무별 맞춤 교육 요청 50%\n"
            "소수 정예 1인 다역 구조에서 코드 기반 자동화 도구 직접 제작 필요  →  과정 B·C·D로 해결"
        ),
        (
            "니즈 3",
            "할루시네이션 검증 및 고품질 결과물 산출",
            "주요 연구원 니즈: 거짓을 가려내고 퀄리티 높은 결과물을 뽑아내는 역량\n"
            "결과물 신뢰성 부족 20%  →  과정 A·C·D 핵심 커리큘럼 반영"
        ),
    ]
    card_h = Inches(1.2)
    gap    = Inches(0.15)
    card_w = Inches(8.89)
    lx     = Inches(0.56)
    left_w = Inches(2.8)   # 배지+제목 영역 너비
    div_x  = lx + left_w
    body_x = div_x + Inches(0.18)
    body_w = card_w - left_w - Inches(0.28)

    for i, (badge, title, body) in enumerate(needs):
        y = Inches(1.19) + i * (card_h + gap)
        rect(sl, lx, y, card_w, card_h, fill=CARD)

        # 배지 (오렌지 태그)
        tb(sl, badge,
           lx + Inches(0.18), y + Inches(0.13),
           Inches(0.75), Inches(0.26),
           size=9, bold=True, color=OR)

        # 제목
        tb(sl, title,
           lx + Inches(0.18), y + Inches(0.44),
           left_w - Inches(0.3), Inches(0.65),
           size=11, bold=True, color=DARK)

        # 세로 구분선
        rect(sl, div_x, y + Inches(0.15),
             Inches(0.005), card_h - Inches(0.3), fill=HEAD)

        # 본문 텍스트 (우측)
        tb(sl, body,
           body_x, y + Inches(0.2),
           body_w, card_h - Inches(0.4),
           size=10, color=DARK)

# ── S07 전체 교육 구조 7회 개요 ───────────────────────────────
def s07(prs):
    sl = blank(prs)
    content_page(sl, "3. 교육 과정 제안", "전체 교육 구조 개요 (7회 / 28시간)",
                 "AI 이해에서 코드 기반 업무 자동화까지 — 단계별 성장 구조")

    # 7회차 표
    headers = ["회차", "과정", "등급", "주제", "대상"]
    rows = [
        ("1회", "과정 A",   "—",   "AI 리터러시 + 프롬프트 실습",        "전사 33명"),
        ("2회", "과정 B-1", "초급", "바이브코딩 환경 구축 + 첫 프로젝트", "전사 33명"),
        ("3회", "과정 B-2", "초급", "데이터 처리 + 업무 자동화 스크립트", "전사 33명"),
        ("4회", "과정 C-1", "중급", "API 연동 + 웹 대시보드",             "실무 담당자"),
        ("5회", "과정 C-2", "중급", "웹앱 구현 + 업무 도구 제작",         "실무 담당자"),
        ("6회", "과정 D-1", "상급", "복합 업무 시스템 + 배포",            "핵심 인력"),
        ("7회", "과정 D-2", "상급", "AI 통합 업무 도구 + PoC 발표",       "핵심 인력"),
    ]
    cw_list = [Inches(0.6), Inches(0.9), Inches(0.7), Inches(4.2), Inches(1.8)]
    xs_list = [Inches(0.56), Inches(1.25), Inches(2.24), Inches(3.03), Inches(7.32)]
    rh = Inches(0.38)

    for j, (hdr, xp) in enumerate(zip(headers, xs_list)):
        rect(sl, xp, Inches(1.19), cw_list[j], rh, fill=HEAD)
        tb(sl, hdr, xp + Inches(0.06), Inches(1.21),
           cw_list[j] - Inches(0.12), rh - Inches(0.08),
           size=9, bold=True, color=DARK)
    for i, row in enumerate(rows):
        ry = Inches(1.57) + i * rh
        for j, (cell, xp) in enumerate(zip(row, xs_list)):
            rect(sl, xp, ry, cw_list[j], rh, fill=DATA)
            tb(sl, cell, xp + Inches(0.06), ry + Inches(0.06),
               cw_list[j] - Inches(0.12), rh - Inches(0.12),
               size=9, color=DARK)

    tb(sl, "총 7회 / 28시간  |  각 4시간  |  월 1회 정기 실습 교육",
       Inches(0.56), Inches(5.2), Inches(8.89), Inches(0.28),
       size=9, color=GRAY)

# ── S08 과정 A ────────────────────────────────────────────────
def s08(prs):
    sl = blank(prs)
    content_page(sl, "3. 교육 과정 제안", "과정 A: AI 리터러시 + 프롬프트 실습",
                 "전사 33명 대상 / 1회 / 4시간")

    info_rows = [
        ("과정명", "생성형 AI 리터러시 & 프롬프트 실전 활용"),
        ("대상",   "전사 임직원 33명 (Lv1~2 입문·초급 55% 대상)"),
        ("시간",   "4시간 이내 (1회)"),
    ]
    lw_col = Inches(1.35)
    vw_col = Inches(7.54)
    lx = Inches(0.56)
    vx = Inches(2.0)
    rh = Inches(0.28)
    for r, (k, v) in enumerate(info_rows):
        ry = Inches(1.19) + r * rh
        rect(sl, lx, ry, lw_col, rh, fill=HEAD)
        tb(sl, k, lx + Inches(0.08), ry + Inches(0.04),
           lw_col - Inches(0.16), rh - Inches(0.08),
           size=9, bold=True, color=DARK)
        rect(sl, vx, ry, vw_col, rh, fill=DATA)
        tb(sl, v, vx + Inches(0.08), ry + Inches(0.04),
           vw_col - Inches(0.16), rh - Inches(0.08),
           size=9, color=DARK)

    base_y = Inches(1.19) + len(info_rows) * rh + Inches(0.2)

    tb(sl, "교육 목표", Inches(0.56), base_y, Inches(3.0), Inches(0.28),
       size=11, bold=True, color=DARK)
    tb(sl, "업무에 바로 활용 가능한 AI 도구 활용 능력을 갖추고, 직무별 프롬프트 설계 역량을 내재화한다. AI의 한계와 원칙을 이해하여 주체적인 AI 활용 문화를 조성한다.",
       Inches(0.56), base_y + Inches(0.3), Inches(8.89), Inches(0.5),
       size=10, color=DARK)

    mid_y = base_y + Inches(0.85)
    half_w = Inches(4.17)

    tb(sl, "커리큘럼 개요", Inches(0.56), mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    curriculum = [
        "• ChatGPT·Claude·Gemini 비교 및 업무 활용 시나리오 이해",
        "• 프롬프트 엔지니어링 원칙 및 직무별 실습 (보고서·이메일·기획안)",
        "• 고급 프롬프트 전략: Chain-of-Thought, Few-shot, Self-consistency ★",
        "• AI 할루시네이션 이해 및 Fact-check 실습 ★ (연구원 핵심 니즈)",
        "• AI 활용 표준화: 조직에서 AI를 안전하고 일관되게 쓰는 가이드라인",
        "• 딥 리서치·Napkin AI·NotebookLM 등 생산성 도구 실습",
    ]
    tb_lines(sl, curriculum,
             Inches(0.56), mid_y + Inches(0.32),
             half_w, Inches(len(curriculum) * 0.3),
             size=9, color=DARK)

    rx = Inches(5.07)
    tb(sl, "기대 효과", rx, mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    effects = [
        ("Lv1~2 → Lv2~3 목표",  "전사 AI 활용 기준선 상향, 개인 역량 편차 해소"),
        ("업무 30%+ 단축",       "문서 작성·요약·기획 등 일상 업무 시간 단축"),
        ("책임 있는 AI 활용",    "AI 맹신 방지 및 올바른 활용 문화 정착"),
    ]
    for i, (eff_title, eff_body) in enumerate(effects):
        ey = mid_y + Inches(0.32) + i * Inches(0.7)
        rect(sl, rx, ey, half_w, Inches(0.62), fill=CARD)
        tb(sl, eff_title, rx + Inches(0.12), ey + Inches(0.07),
           half_w - Inches(0.24), Inches(0.26), size=10, bold=True, color=OR)
        tb(sl, eff_body, rx + Inches(0.12), ey + Inches(0.35),
           half_w - Inches(0.24), Inches(0.24), size=9, color=DARK)

# ── S09 과정 B: 바이브코딩 초급 ──────────────────────────────
def s09(prs):
    sl = blank(prs)
    content_page(sl, "3. 교육 과정 제안", "과정 B: 바이브코딩 초급 (2회)",
                 "전사 33명 대상 / 2회 / 각 4시간")

    info_rows = [
        ("과정명", "바이브코딩으로 시작하는 업무 자동화"),
        ("대상",   "전사 임직원 33명 (Lv2~3 역량 기반 실습)"),
        ("시간",   "4시간 이내 × 2회"),
    ]
    lw_col = Inches(1.35)
    vw_col = Inches(7.54)
    lx = Inches(0.56)
    vx = Inches(2.0)
    rh = Inches(0.28)
    for r, (k, v) in enumerate(info_rows):
        ry = Inches(1.19) + r * rh
        rect(sl, lx, ry, lw_col, rh, fill=HEAD)
        tb(sl, k, lx + Inches(0.08), ry + Inches(0.04),
           lw_col - Inches(0.16), rh - Inches(0.08),
           size=9, bold=True, color=DARK)
        rect(sl, vx, ry, vw_col, rh, fill=DATA)
        tb(sl, v, vx + Inches(0.08), ry + Inches(0.04),
           vw_col - Inches(0.16), rh - Inches(0.08),
           size=9, color=DARK)

    mid_y = Inches(1.19) + len(info_rows) * rh + Inches(0.2)
    half_w = Inches(4.17)

    # 1회차
    tb(sl, "1회차: 환경 구축 + 첫 프로젝트", Inches(0.56), mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    b1_items = [
        "• 바이브코딩 개념 이해 (AI와 대화로 프로그래밍)",
        "• VS Code + Claude Code 설치 및 설정",
        "• PRD/BRD 작성법: 자연어로 요구사항 구조화",
        "• 실습: 간단한 업무 도구 만들기 (계산기, 할 일 목록)",
        "• 컨트로맥스 실습: ServoMax 제품 스펙 비교표 생성기",
    ]
    tb_lines(sl, b1_items,
             Inches(0.56), mid_y + Inches(0.32),
             half_w, Inches(len(b1_items) * 0.3),
             size=9, color=DARK)

    # 2회차
    rx = Inches(5.07)
    tb(sl, "2회차: 데이터 처리 + 업무 자동화", rx, mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    b2_items = [
        "• CSV/엑셀 데이터 읽기·분석·변환 자동화",
        "• 파일 정리·이름 변경·백업 자동화 스크립트",
        "• 이메일·보고서 템플릿 자동 생성",
        "• 실습: 본인 반복 업무 자동화 프로젝트 구현 및 발표",
        "• 컨트로맥스 실습: 영문 이메일 자동 작성기 / 회의록 정리",
    ]
    tb_lines(sl, b2_items,
             rx, mid_y + Inches(0.32),
             half_w, Inches(len(b2_items) * 0.3),
             size=9, color=DARK)

    tb(sl, "목표: Lv2~3 → Lv3  |  산출물: 개인별 업무 자동화 스크립트",
       Inches(0.56), Inches(5.2), Inches(8.89), Inches(0.28),
       size=9, color=GRAY)

# ── S10 과정 C: 바이브코딩 중급 ──────────────────────────────
def s10(prs):
    sl = blank(prs)
    content_page(sl, "3. 교육 과정 제안", "과정 C: 바이브코딩 중급 (2회)",
                 "실무 담당자 대상 / 2회 / 각 4시간")

    info_rows = [
        ("과정명", "바이브코딩으로 만드는 웹 대시보드 & 업무 도구"),
        ("대상",   "실무 담당자 (과정 B 수료 또는 동등 수준) — Lv3 역량 대상"),
        ("시간",   "4시간 이내 × 2회"),
    ]
    lw_col = Inches(1.35)
    vw_col = Inches(7.54)
    lx = Inches(0.56)
    vx = Inches(2.0)
    rh = Inches(0.28)
    for r, (k, v) in enumerate(info_rows):
        ry = Inches(1.19) + r * rh
        rect(sl, lx, ry, lw_col, rh, fill=HEAD)
        tb(sl, k, lx + Inches(0.08), ry + Inches(0.04),
           lw_col - Inches(0.16), rh - Inches(0.08),
           size=9, bold=True, color=DARK)
        rect(sl, vx, ry, vw_col, rh, fill=DATA)
        tb(sl, v, vx + Inches(0.08), ry + Inches(0.04),
           vw_col - Inches(0.16), rh - Inches(0.08),
           size=9, color=DARK)

    mid_y = Inches(1.19) + len(info_rows) * rh + Inches(0.2)
    half_w = Inches(4.17)

    tb(sl, "1회차: API 연동 + 웹 대시보드", Inches(0.56), mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    c1_items = [
        "• LLM API 기초 (Claude/ChatGPT API 호출)",
        "• Streamlit 대시보드 구현",
        "• 업무 데이터 시각화 + 자동 보고서 생성",
        "• 실습: 팀 업무 데이터 대시보드 구현",
        "• 컨트로맥스 실습: 생산 일정 관리 캘린더 웹앱",
    ]
    tb_lines(sl, c1_items,
             Inches(0.56), mid_y + Inches(0.32),
             half_w, Inches(len(c1_items) * 0.3),
             size=9, color=DARK)

    rx = Inches(5.07)
    tb(sl, "2회차: 웹앱 + 업무 도구 제작", rx, mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    c2_items = [
        "• 풀스택 웹앱 구현 (바이브코딩으로 프론트+백엔드)",
        "• 할루시네이션 검증 시스템 구현 (교차 검증) ★",
        "• 실습: 팀별 업무 웹앱 제작 및 시연",
        "• 컨트로맥스 실습: 고객 문의 자동 응답 도구",
        "• 컨트로맥스 실습: ServoMax 제품 소개 웹페이지",
    ]
    tb_lines(sl, c2_items,
             rx, mid_y + Inches(0.32),
             half_w, Inches(len(c2_items) * 0.3),
             size=9, color=DARK)

    tb(sl, "목표: Lv3 → Lv3~4  |  산출물: 팀별 웹 대시보드·업무 웹앱",
       Inches(0.56), Inches(5.2), Inches(8.89), Inches(0.28),
       size=9, color=GRAY)

# ── S11 과정 D: 바이브코딩 상급 ──────────────────────────────
def s11(prs):
    sl = blank(prs)
    content_page(sl, "3. 교육 과정 제안", "과정 D: 바이브코딩 상급 (2회)",
                 "핵심 인력 대상 / 2회 / 각 4시간")

    info_rows = [
        ("과정명", "바이브코딩으로 구축하는 AI 통합 업무 시스템"),
        ("대상",   "핵심 인력 / DX 추진 담당자 (과정 C 수료 또는 동등) — Lv3~4"),
        ("시간",   "4시간 이내 × 2회"),
    ]
    lw_col = Inches(1.35)
    vw_col = Inches(7.54)
    lx = Inches(0.56)
    vx = Inches(2.0)
    rh = Inches(0.28)
    for r, (k, v) in enumerate(info_rows):
        ry = Inches(1.19) + r * rh
        rect(sl, lx, ry, lw_col, rh, fill=HEAD)
        tb(sl, k, lx + Inches(0.08), ry + Inches(0.04),
           lw_col - Inches(0.16), rh - Inches(0.08),
           size=9, bold=True, color=DARK)
        rect(sl, vx, ry, vw_col, rh, fill=DATA)
        tb(sl, v, vx + Inches(0.08), ry + Inches(0.04),
           vw_col - Inches(0.16), rh - Inches(0.08),
           size=9, color=DARK)

    mid_y = Inches(1.19) + len(info_rows) * rh + Inches(0.2)
    half_w = Inches(4.17)

    tb(sl, "1회차: 복합 업무 시스템 + 배포", Inches(0.56), mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    d1_items = [
        "• 다중 기능 웹앱 설계 (인증, DB, API 통합)",
        "• 데이터베이스 연동 (Supabase 등)",
        "• 배포 실습 (Vercel / Streamlit Cloud)",
        "• 실습: 팀 업무 관리 시스템 구현 및 실제 배포",
        "• 컨트로맥스 실습: 고객 주문·납품 추적 시스템",
    ]
    tb_lines(sl, d1_items,
             Inches(0.56), mid_y + Inches(0.32),
             half_w, Inches(len(d1_items) * 0.3),
             size=9, color=DARK)

    rx = Inches(5.07)
    tb(sl, "2회차: AI 통합 업무 도구 + PoC 발표", rx, mid_y, half_w, Inches(0.28),
       size=11, bold=True, color=DARK)
    d2_items = [
        "• AI API 활용 사내 챗봇·문서 도우미 구현",
        "• 복합 자동화 시스템 설계 (다수 API 연동)",
        "• 팀별 최종 프로젝트 완성",
        "• PoC 발표 + 상호 피드백 + 최종 역량 평가",
        "• 컨트로맥스 실습: 사내 기술 문서 검색 챗봇",
    ]
    tb_lines(sl, d2_items,
             rx, mid_y + Inches(0.32),
             half_w, Inches(len(d2_items) * 0.3),
             size=9, color=DARK)

    tb(sl, "목표: Lv3~4 → Lv4~5  |  산출물: AI 업무 도구 + PoC + 역량 평가 보고서",
       Inches(0.56), Inches(5.2), Inches(8.89), Inches(0.28),
       size=9, color=GRAY)

# ── S12 역량 평가 구조 ────────────────────────────────────────
def s12(prs):
    sl = blank(prs)
    content_page(
        sl, "4. AI 역량 평가 및 인사평가 연계", "AI 역량 평가 — 사전·사후 구조",
        "모든 과정에 사전 진단 → 교육 → 사후 평가 구조를 적용합니다. 클라이언트 요청사항 ③④를 직접 충족합니다."
    )

    steps = [
        ("사전 진단",   "레벨 측정\n(설문 + 실기)"),
        ("교육 실시",   "맞춤형 과정\nA / B / C / D"),
        ("사후 평가",   "레벨 재측정\n(설문 + 실기)"),
        ("인사평가 연계", "성과 지표 보고\n(레벨업 기준)"),
    ]
    bw = Inches(2.02)
    bh = Inches(1.1)
    gap = Inches(0.27)
    sx = Inches(0.56)
    for i, (title, body) in enumerate(steps):
        x = sx + i * (bw + gap)
        y = Inches(1.32)
        rect(sl, x, y, bw, bh, fill=CARD)
        tb(sl, title, x + Inches(0.1), y + Inches(0.1),
           bw - Inches(0.2), Inches(0.28),
           size=10, bold=True, color=OR, align=PP_ALIGN.CENTER)
        tb(sl, body, x + Inches(0.1), y + Inches(0.44),
           bw - Inches(0.2), Inches(0.55),
           size=9, color=DARK, align=PP_ALIGN.CENTER)
        if i < 3:
            tb(sl, "→",
               x + bw + Inches(0.08), y + bh / 2 - Inches(0.14),
               Inches(0.12), Inches(0.28),
               size=12, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    left_y = Inches(2.7)
    tb(sl, "평가 방법", Inches(0.56), left_y, Inches(4.1), Inches(0.28),
       size=11, bold=True, color=DARK)
    eval_items = [
        "• 사전: 진단 설문 + 실기 과제 (AI 도구 활용 수행 평가)",
        "• 사후: 동일 기준 재평가 + 레벨업 인정",
        "• 워카 발급 AI 역량 수료 증명서 (레벨 명시)",
    ]
    tb_lines(sl, eval_items,
             Inches(0.56), left_y + Inches(0.34),
             Inches(4.1), Inches(len(eval_items) * 0.38),
             size=10, color=DARK)

    rx = Inches(5.1)
    tb(sl, "인사평가 연계 방안", rx, left_y, Inches(4.3), Inches(0.28),
       size=11, bold=True, color=DARK)
    hr_items = [
        "• 레벨업 1단계 달성 → 역량 개발 목표 달성 ✓",
        "• 레벨업 2단계 이상 → 우수 역량 개발 ✓",
        "• Lv4 도달 시 → 사내 AI 챔피언 후보군 등재",
        "• 구체적 기준: 귀사 인사 정책 협의 조정 가능",
    ]
    tb_lines(sl, hr_items,
             rx, left_y + Inches(0.34),
             Inches(4.3), Inches(len(hr_items) * 0.38),
             size=10, color=DARK)

# ── S13 역량 레벨 정의 ───────────────────────────────────────
def s13(prs):
    sl = blank(prs)
    content_page(sl, "4. AI 역량 평가 및 인사평가 연계", "AI 역량 레벨 정의 (Lv1 ~ Lv5)")

    level_data = [
        ["레벨", "명칭",  "정의",                 "주요 역량"],
        ["Lv1", "입문",  "AI 도구 인지·기본 사용", "ChatGPT 기본 질문, 단순 요청"],
        ["Lv2", "초급",  "직무 활용 시작",          "프롬프트 작성, 업무 문서 생성"],
        ["Lv3", "중급",  "자립적 AI 활용",          "프롬프트 최적화, 반복 업무 자동화"],
        ["Lv4", "고급",  "팀 내 전파 가능",         "웹앱 제작, API 연동, 결과물 검증"],
        ["Lv5", "전문가", "조직 AI 내재화 주도",     "AI 전략 수립, 커스텀 솔루션 구현·배포"],
    ]
    cw = [Inches(0.8), Inches(0.9), Inches(2.8), Inches(4.38)]
    xs = [Inches(0.56), Inches(1.45), Inches(2.44), Inches(5.33)]
    rh = Inches(0.42)

    for r, row in enumerate(level_data):
        ry = Inches(1.19) + r * rh
        for j, (cell, xp) in enumerate(zip(row, xs)):
            bg = HEAD if r == 0 else DATA
            rect(sl, xp, ry, cw[j], rh, fill=bg)
            tb(sl, cell, xp + Inches(0.08), ry + Inches(0.07),
               cw[j] - Inches(0.16), rh - Inches(0.14),
               size=10, bold=(r == 0), color=DARK)

    tb(sl, "교육 과정별 목표 레벨",
       Inches(0.56), Inches(3.82), Inches(8.89), Inches(0.28),
       size=11, bold=True, color=DARK)

    goal_data = [
        ["과정",   "입과 기준", "수료 목표",  "비고"],
        ["과정 A", "Lv1~2",    "→ Lv2~3", "전사 33명 (AI 리터러시)"],
        ["과정 B", "Lv2~3",    "→ Lv3",   "전사 33명 (바이브코딩 초급)"],
        ["과정 C", "Lv3",      "→ Lv3~4", "실무 담당자 (바이브코딩 중급)"],
        ["과정 D", "Lv3~4",    "→ Lv4~5", "핵심 인력 (바이브코딩 상급)"],
    ]
    gcw = [Inches(1.2), Inches(1.2), Inches(1.2), Inches(4.78)]
    gxs = [Inches(0.56), Inches(1.85), Inches(3.14), Inches(4.43)]
    grh = Inches(0.3)

    for r, row in enumerate(goal_data):
        ry = Inches(4.18) + r * grh
        for j, (cell, xp) in enumerate(zip(row, gxs)):
            bg = HEAD if r == 0 else DATA
            rect(sl, xp, ry, gcw[j], grh, fill=bg)
            tb(sl, cell, xp + Inches(0.08), ry + Inches(0.04),
               gcw[j] - Inches(0.16), grh - Inches(0.08),
               size=9, bold=(r == 0), color=DARK)

# ── S14 인사평가 연계 ─────────────────────────────────────────
def s14(prs):
    sl = blank(prs)
    content_page(sl, "4. AI 역량 평가 및 인사평가 연계", "인사평가 연계 방안",
                 "AI 역량 평가 결과를 인사평가에 공식 반영하는 체계입니다.")

    rows = [
        ("레벨업 1단계 달성",     "역량 개발 목표 달성 ✓"),
        ("레벨업 2단계 이상",     "우수 역량 개발 ✓"),
        ("Lv4 도달",             "사내 AI 챔피언 후보군 등재 / 팀 내 전파자 역할 공식화"),
        ("수료 미달 (사유 있음)", "별도 보완 교육 권고 (귀사 인사 정책 협의)"),
    ]
    for i, (cond, result) in enumerate(rows):
        iy = Inches(1.32) + i * Inches(0.88)
        rect(sl, Inches(0.56), iy, Inches(2.8), Inches(0.7), fill=HEAD)
        tb(sl, cond, Inches(0.65), iy + Inches(0.18),
           Inches(2.62), Inches(0.38),
           size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        tb(sl, "→", Inches(3.5), iy + Inches(0.2),
           Inches(0.4), Inches(0.35),
           size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        rect(sl, Inches(4.0), iy, Inches(5.44), Inches(0.7), fill=CARD)
        tb(sl, result, Inches(4.12), iy + Inches(0.16),
           Inches(5.2), Inches(0.4), size=10, color=DARK)

    tb(sl, "※ 구체적인 반영 기준은 귀사 인사 정책에 맞게 협의 조정 가능합니다.",
       Inches(0.56), Inches(5.0), Inches(8.89), Inches(0.3),
       size=9, color=GRAY)

# ── S15 실행 로드맵 ───────────────────────────────────────────
def s15(prs):
    sl = blank(prs)
    content_page(sl, "5. 실행 로드맵", "실행 로드맵",
                 "선호 교육 형태: 월 1회, 4시간 이내 정기 실습 교육")

    phases = [
        {
            "phase": "Phase 1",
            "period": "기초 역량 확보",
            "items": [
                "과정 A 전사 실시 (33명 전원, 4시간 × 1회)",
                "사전 역량 진단 (Lv 측정)",
                "수료 후 사내 AI 활용 가이드라인 공동 수립",
            ],
            "output": "AI 활용 가이드라인\n개인 프롬프트 라이브러리",
        },
        {
            "phase": "Phase 2",
            "period": "바이브코딩 초급 실습",
            "items": [
                "과정 B 전사 대상 (4시간 × 2회)",
                "1회차: 환경 구축 + 첫 프로젝트",
                "2회차: 데이터 처리 + 업무 자동화",
                "중간 역량 평가 (레벨업 현황 확인)",
            ],
            "output": "개인별 업무 자동화 스크립트",
        },
        {
            "phase": "Phase 3",
            "period": "바이브코딩 중급 실습",
            "items": [
                "과정 C 실무 담당자 대상 (4시간 × 2회)",
                "1회차: API 연동 + 웹 대시보드",
                "2회차: 웹앱 + 업무 도구 제작",
            ],
            "output": "팀별 웹 대시보드·업무 웹앱",
        },
        {
            "phase": "Phase 4",
            "period": "바이브코딩 상급 + PoC",
            "items": [
                "과정 D 핵심 인력 대상 (4시간 × 2회)",
                "1회차: 복합 시스템 + 배포",
                "2회차: AI 도구 + PoC 발표",
                "사후 역량 평가 + 인사평가 연계 보고서",
            ],
            "output": "AI 업무 도구 + PoC\n역량 평가 보고서",
        },
    ]
    # 4 Phase: 2x2 배치
    cw = Inches(4.3)
    ch = Inches(2.0)
    gap_x = Inches(0.29)
    gap_y = Inches(0.2)
    positions = [
        (Inches(0.56), Inches(1.19)),
        (Inches(0.56) + cw + gap_x, Inches(1.19)),
        (Inches(0.56), Inches(1.19) + ch + gap_y),
        (Inches(0.56) + cw + gap_x, Inches(1.19) + ch + gap_y),
    ]
    for i, p in enumerate(phases):
        x, y = positions[i]
        rect(sl, x, y, cw, ch, fill=CARD)
        tb(sl, p["phase"], x + Inches(0.15), y + Inches(0.1),
           cw - Inches(0.3), Inches(0.26), size=11, bold=True, color=OR)
        tb(sl, p["period"], x + Inches(0.15), y + Inches(0.38),
           cw - Inches(0.3), Inches(0.22), size=9, color=GRAY)
        bullet_items = [f"• {item}" for item in p["items"]]
        tb_lines(sl, bullet_items,
                 x + Inches(0.15), y + Inches(0.64),
                 cw - Inches(0.3), Inches(len(bullet_items) * 0.28),
                 size=8, color=DARK)
        tb(sl, "산출물: " + p["output"].replace("\n", " / "),
           x + Inches(0.15), y + ch - Inches(0.28),
           cw - Inches(0.3), Inches(0.24), size=8, bold=True, color=OR)

# ── S16 워카 소개 ─────────────────────────────────────────────
def s16(prs):
    sl = blank(prs)
    content_page(sl, "6. 워카 소개 & 레퍼런스", "워카 소개")

    stats = [
        ("12,000명+", "누적 수강생"),
        ("11,000시간+", "누적 교육 시간"),
        ("75개+",    "교육 기관 수"),
        ("2022",     "설립"),
    ]
    cw = Inches(2.02)
    gap = Inches(0.27)
    for i, (val, label) in enumerate(stats):
        x = Inches(0.56) + i * (cw + gap)
        y = Inches(1.19)
        rect(sl, x, y, cw, Inches(0.96), fill=CARD)
        tb(sl, val, x, y + Inches(0.06), cw, Inches(0.52),
           size=20, bold=True, color=OR, align=PP_ALIGN.CENTER)
        tb(sl, label, x, y + Inches(0.6), cw, Inches(0.32),
           size=9, color=DARK, align=PP_ALIGN.CENTER)

    tb(sl,
       "워카는 실무에 바로 적용할 수 있는 AI·데이터 교육을 제공하는 교육 전문 기업입니다.\n"
       "삼성전자, POSCO인터내셔널, SK하이닉스, 하나금융그룹, 현대위아 등 국내 대표 기업과 공공기관에서 검증된 교육 솔루션을 운영합니다.",
       Inches(0.56), Inches(2.28), Inches(8.89), Inches(0.65),
       size=10, color=DARK)

    tb(sl, "주요 고객사",
       Inches(0.56), Inches(3.06), Inches(8.89), Inches(0.28),
       size=11, bold=True, color=DARK)

    clients = ["삼성전자", "SK하이닉스", "POSCO인터내셔널", "하나금융그룹",
               "현대위아", "차의과대학교", "공공기관 다수"]
    ccw = Inches(1.18)
    cgap = Inches(0.1)
    for i, c in enumerate(clients):
        cx = Inches(0.56) + i * (ccw + cgap)
        rect(sl, cx, Inches(3.42), ccw, Inches(0.36), fill=DATA)
        tb(sl, c, cx, Inches(3.46), ccw, Inches(0.28),
           size=9, color=DARK, align=PP_ALIGN.CENTER)

# ── S17 유사 교육 사례 ────────────────────────────────────────
def s17(prs):
    sl = blank(prs)
    content_page(sl, "6. 워카 소개 & 레퍼런스", "유사 교육 사례")

    cases = [
        {
            "title": "SK하이닉스",
            "sub":   "GPT 활용 업무 능력 향상 (2026)",
            "body":  (
                "대상: 신입사원 180명 / 5시간  |  프롬프트 엔지니어링  |  OA 자동화 실습\n"
                "결과: 신입사원 AI 리터러시 및 실전 활용 역량 확보"
            ),
        },
        {
            "title": "포스코인터내셔널",
            "sub":   "바이브코딩 업무 자동화 (2026)",
            "body":  (
                "대상: 비개발 신입사원 20명 / 7시간 × 2회차  |  VS Code + Claude Code  |  업무 자동화 실습\n"
                "결과: 비개발자도 메일 요약·보고서 자동화·웹앱 구현 달성"
            ),
        },
        {
            "title": "차의과대학교",
            "sub":   "AI 리터러시 및 직무 역량 교육 (2026)",
            "body":  (
                "대상: 전 직원·임원 80명 / 4시간 × 6회차 (난이도별)  |  프롬프트 엔지니어링  |  GPTs 제작  |  보고서 자동화\n"
                "결과: 전 직원 AI 역량 내재화 / 임원 대상 별도 커리큘럼 운영"
            ),
        },
    ]
    card_h = Inches(1.2)
    gap    = Inches(0.15)
    card_w = Inches(8.89)
    lx     = Inches(0.56)
    left_w = Inches(2.4)   # 기관명+부제 영역 너비
    div_x  = lx + left_w
    body_x = div_x + Inches(0.18)
    body_w = card_w - left_w - Inches(0.28)

    for i, c in enumerate(cases):
        y = Inches(1.19) + i * (card_h + gap)
        rect(sl, lx, y, card_w, card_h, fill=CARD)

        # 기관명
        tb(sl, c["title"],
           lx + Inches(0.18), y + Inches(0.17),
           left_w - Inches(0.3), Inches(0.36),
           size=13, bold=True, color=DARK)

        # 부제 (회색 소문자)
        tb(sl, c["sub"],
           lx + Inches(0.18), y + Inches(0.60),
           left_w - Inches(0.3), Inches(0.44),
           size=9, color=GRAY)

        # 세로 구분선
        rect(sl, div_x, y + Inches(0.15),
             Inches(0.005), card_h - Inches(0.3), fill=HEAD)

        # 본문 텍스트 (우측)
        tb(sl, c["body"],
           body_x, y + Inches(0.27),
           body_w, card_h - Inches(0.44),
           size=9, color=DARK)

# ── S18 교육 성과 사례 ────────────────────────────────────────
def s18_perf(prs):
    sl = blank(prs)
    content_page(sl, "6. 워카 소개 & 레퍼런스", "교육을 통한 업무 생산성 향상 사례")

    cases = [
        {
            "title": "S기업",
            "sub":   "폐수 처리 약품 용량 예측",
            "body":  "AI 기반 예측 모델 직접 구현 — 공장 폐수 처리 약품 투입량 예측 자동화",
            "perf":  "예측 모델 구현 → 00억 비용절감",
        },
        {
            "title": "H기업",
            "sub":   "트럼프 관세 대응 자료조사",
            "body":  "AI 딥 리서치 + 바이브코딩으로 자료조사·보고서·발표 자료 일괄 제작",
            "perf":  "기존보다 4배 빠른 자료조사·보고서·발표 제작",
        },
        {
            "title": "L기업",
            "sub":   "공장 출입 명단 입력 자동화",
            "body":  "크롤링 자동화 스크립트 제작 — 매일 반복되는 수작업 출입 명단 등록 자동 처리",
            "perf":  "매일 1시간 업무 시간 절약",
        },
    ]
    card_h = Inches(1.2)
    gap    = Inches(0.15)
    card_w = Inches(8.89)
    lx     = Inches(0.56)
    left_w = Inches(2.4)
    div_x  = lx + left_w
    body_x = div_x + Inches(0.18)
    body_w = card_w - left_w - Inches(0.28)

    for i, c in enumerate(cases):
        y = Inches(1.19) + i * (card_h + gap)
        rect(sl, lx, y, card_w, card_h, fill=CARD)

        # 기업 이니셜 (볼드)
        tb(sl, c["title"],
           lx + Inches(0.18), y + Inches(0.17),
           left_w - Inches(0.3), Inches(0.36),
           size=13, bold=True, color=DARK)

        # 과제 요약 (회색)
        tb(sl, c["sub"],
           lx + Inches(0.18), y + Inches(0.60),
           left_w - Inches(0.3), Inches(0.44),
           size=9, color=GRAY)

        # 세로 구분선
        rect(sl, div_x, y + Inches(0.15),
             Inches(0.005), card_h - Inches(0.3), fill=HEAD)

        # 상세 내용 (우측 상단)
        tb(sl, c["body"],
           body_x, y + Inches(0.17),
           body_w, Inches(0.46),
           size=9, color=DARK)

        # 성과 하이라이트 (오렌지 볼드)
        tb(sl, c["perf"],
           body_x, y + Inches(0.70),
           body_w, Inches(0.38),
           size=10, bold=True, color=OR)


# ── S19 담당자 연락처 ─────────────────────────────────────────
def s19(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=DARK)

    tb(sl, "읽어주셔서 감사합니다.",
       Inches(0.56), Inches(2.32), Inches(7.82), Inches(0.65),
       size=26, bold=True, color=WH)

    rows = [
        ("담당자",  "교육팀 채진백 팀장"),
        ("연락처",  "010-2326-4348  |  jinbaek@woka.kr"),
        ("홈페이지", "woka.kr"),
    ]
    for i, (k, v) in enumerate(rows):
        y = Inches(3.18) + i * Inches(0.46)
        tb(sl, k, Inches(0.56), y, Inches(1.2), Inches(0.36),
           size=12, color=GRAY)
        tb(sl, v, Inches(1.9), y, Inches(7.0), Inches(0.36),
           size=12, bold=True, color=WH)

    tb(sl, "주식회사 워카 (Woka)  ·  woka.kr",
       Inches(0.56), Inches(5.22), Inches(8.89), Inches(0.3),
       size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── main ──────────────────────────────────────────────────────
def main():
    prs = prs_new()
    funcs = [
        (s01, "표지"),
        (s02, "목차"),
        (s03, "교육 의뢰 배경"),
        (s04, "설문 핵심 수치"),
        (s05, "업무 영역 & 장벽"),
        (s06, "교육 니즈 3가지"),
        (s07, "전체 교육 구조 7회"),
        (s08, "과정 A: AI 리터러시"),
        (s09, "과정 B: 바이브코딩 초급"),
        (s10, "과정 C: 바이브코딩 중급"),
        (s11, "과정 D: 바이브코딩 상급"),
        (s12, "역량 평가 구조"),
        (s13, "레벨 정의"),
        (s14, "인사평가 연계"),
        (s15, "실행 로드맵"),
        (s16, "워카 소개"),
        (s17, "유사 사례"),
        (s18_perf, "교육 성과 사례"),
        (s19, "담당자 연락처"),
    ]
    for i, (fn, label) in enumerate(funcs, 1):
        fn(prs)
        print(f"  {i:2d}/19  {label}")

    out = "/Users/chaejinbaek/Desktop/바이브코딩/edu_proposal_v2/교육_제안서_2안.pptx"
    prs.save(out)
    print(f"\n✅ 저장 완료: {out}")
    print(f"   슬라이드 수: {len(prs.slides)}")
    print(f"   슬라이드 크기: {prs.slide_width.inches:.2f}in × {prs.slide_height.inches:.2f}in")

if __name__ == "__main__":
    main()
